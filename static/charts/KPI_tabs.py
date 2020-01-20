# -*- coding: utf-8 -*-
"""
Created on Tue Jan  1 13:51:16 2019
@author: neel_
"""

import numpy as np
import pandas as pd
import os
import time
#import datetime
import operator
import pandas as pd
import numpy as np
from scipy import stats
import matplotlib.pyplot as plt
#%matplotlib inline
import math
import seaborn as sns
#import catplot
sns.set(style="ticks")
from datetime import datetime
import matplotlib.mlab as mlab
import pylab as pl
from scipy.stats import norm 
from datetime import timedelta
import csv
import openpyxl
from openpyxl import load_workbook
import pptx
import pptx.util
import glob
import scipy.misc



#Set working directory
# aPath='/home/dhananjai/Desktop/TryPyscript'
aPath='/var/django/StagingData/static/charts/'
os.chdir(aPath)


#Create erroneous session list, Reading data generated at validation stage
#dfer=pd.read_csv('Error_data.csv')
#err_sess=dfer.groupby(['ses_id']).agg({'DeviceId':'nunique'}).reset_index()
#
##Session ID with erroneous data, to be excluded from final analysis
#ErSess_list=err_sess['ses_id'].tolist()

#Read data with unique session id, Reading data generated at validation stage
data_process=pd.read_csv('Clean_data.csv')
df_music=pd.read_csv('CleanData_WithDur.csv')


#Read the session wise data
# df=pd.read_csv('Data_sessionID.csv')#----------------------------------
#Replace station MHz in quotes to without quotes
data_process.station=data_process.station.str.replace('"','')
#
# df=df[df['ses_id'].isin(ErSess_list)==False]#----------------------------------
# list(df)#----------------------------------
# data_process=df.copy()#----------------------------------
data_process['station']=data_process['station'].astype(str)
data_process['station_key']=data_process['city']+data_process['station']

d_station=pd.read_excel('Input_DataValues.xlsx', sheet_name='Radio')
d_station['Station']=d_station['Station'].astype(str)
d_station['station_key']=d_station['segment_City']+d_station['Station']
data_process=pd.merge(data_process,d_station,how='left',on=['station_key'])

list(data_process)
temp=data_process[data_process['keyname']=='"FM_Tuned"']
temp['Station Name']
#Create tables
d_mop=pd.read_excel('Input_DataValues.xlsx', sheet_name='Mob_op')
data_process=pd.merge(data_process,d_mop,how='left',on=['mobileoperator'])
data_process.drop(columns =["Unnamed: 2",'Unnamed: 3','mobileoperator.1','Texts in raw data'], inplace = True)

d_source=pd.read_excel('Input_DataValues.xlsx', sheet_name='source')
data_process=pd.merge(data_process,d_source,how='left',on=['source'])

d_musapp=pd.read_excel('Input_DataValues.xlsx', sheet_name='app')
data_process=pd.merge(data_process,d_musapp,how='left',on=['app'])


d_make=pd.read_excel('Input_DataValues.xlsx', sheet_name='Dev_make')
data_process=pd.merge(data_process,d_make,how='left',on=['devicemake'])
#data_process.to_csv('test.csv',index=False)


data_process["all"]='All device users'


list(data_process)

temp=data_process.groupby(['keyname','created_date']).agg({'deviceid':'nunique'}).reset_index()
temp.rename(columns={'keyname': 'Event wise #Device users'}, inplace=True)
temp1=data_process.groupby(['all','created_date']).agg({'deviceid':'nunique'}).reset_index()
temp1.rename(columns={'all': 'Event wise #Device users'}, inplace=True)

ctab1=temp.pivot_table(index=['Event wise #Device users'],columns=['created_date'], values='deviceid')
ctab1t=temp1.pivot_table(index=['Event wise #Device users'],columns=['created_date'], values='deviceid')
ctab1=ctab1.append(ctab1t)
                              
temp=data_process.groupby(['mob_op_report','created_date']).agg({'deviceid':'nunique'}).reset_index()
temp.rename(columns={'mob_op_report': 'Mob_operator wise #Device users'}, inplace=True)
temp1=data_process.groupby(['all','created_date']).agg({'deviceid':'nunique'}).reset_index()
temp1.rename(columns={'all': 'Mob_operator wise #Device users'}, inplace=True)

ctab2=temp.pivot_table(index=['Mob_operator wise #Device users'],columns=['created_date'], values='deviceid')
ctab2t=temp1.pivot_table(index=['Mob_operator wise #Device users'],columns=['created_date'], values='deviceid')
ctab2=ctab2.append(ctab2t)


temp=data_process.groupby(['Sourc_rep','created_date']).agg({'deviceid':'nunique'}).reset_index()
temp.rename(columns={'Sourc_rep': 'Mode_listen - #Device users(Among who listened FM/Audio)'}, inplace=True)
temp1=data_process[data_process['keyname'].isin(['"FM_Tuned"','"Audio_Tuned"'])==True]
temp1=temp1.groupby(['all','created_date']).agg({'deviceid':'nunique'}).reset_index()
temp1.rename(columns={'all': 'Mode_listen - #Device users(Among who listened FM/Audio)'}, inplace=True)

ctab3=temp.pivot_table(index=['Mode_listen - #Device users(Among who listened FM/Audio)'],columns=['created_date'], values='deviceid')
ctab3t=temp1.pivot_table(index=['Mode_listen - #Device users(Among who listened FM/Audio)'],columns=['created_date'], values='deviceid')
ctab3=ctab3.append(ctab3t)

#Album & Song
temp=data_process.groupby(['album','created_date']).agg({'deviceid':'nunique'}).reset_index()
temp.rename(columns={'album': 'Album - #Device users'}, inplace=True)

ctab4=temp.pivot_table(index=['Album - #Device users'],columns=['created_date'], values='deviceid')

song_excl=pd.read_excel('Input_DataValues.xlsx', sheet_name='Song_excl')
sng_exclL=song_excl['Song'].tolist()
temp=data_process[data_process['song'].str.contains('/storage/emulated')==False]                  
temp=temp[temp['song'].isin(sng_exclL)==False]
temp=temp.groupby(['song','created_date']).agg({'deviceid':'nunique'}).reset_index()
temp.rename(columns={'song': 'Song - #Device users'}, inplace=True)

ctab5=temp.pivot_table(index=['Song - #Device users'],columns=['created_date'], values='deviceid')

#Music streaming app
list(data_process)
temp=data_process[data_process['keyname']=='"Audio_Tuned"']
#temp['Music_app']=temp["app"].copy()
#temp.loc[temp["app"].str.contains("Jio|Savvn|savvn|Saavn"), "Music_app"] = "Jio/Savvn"
temp=temp.groupby(['app_rep','created_date']).agg({'deviceid':'nunique'}).reset_index()
temp.rename(columns={'app_rep': 'Music_app - #Device users(Tuned audio)'}, inplace=True)

temp1=data_process[data_process['keyname']=='"Audio_Tuned"']
#temp1['Music_app']=temp1["app"].copy()
temp1=temp1.groupby(['all','created_date']).agg({'deviceid':'nunique'}).reset_index()
temp1.rename(columns={'all': 'Music_app - #Device users(Tuned audio)'}, inplace=True)

ctab6t1=temp.pivot_table(index=['Music_app - #Device users(Tuned audio)'],columns=['created_date'], values='deviceid')
ctab6t2=temp1.pivot_table(index=['Music_app - #Device users(Tuned audio)'],columns=['created_date'], values='deviceid')
ctab6=ctab6t1.append(ctab6t2)

#FM/Radio users
temp=data_process[data_process['keyname']=='"FM_Tuned"']
temp=temp.groupby(['Station Name','created_date']).agg({'deviceid':'nunique'}).reset_index()
temp.rename(columns={'Station Name': 'Radio_station - #Device users(Tuned FM/Radio)'}, inplace=True)

temp1=data_process[data_process['keyname']=='"FM_Tuned"']
temp1=temp1.groupby(['all','created_date']).agg({'deviceid':'nunique'}).reset_index()
temp1.rename(columns={'all': 'Radio_station - #Device users(Tuned FM/Radio)'}, inplace=True)

ctab7t1=temp.pivot_table(index=['Radio_station - #Device users(Tuned FM/Radio)'],columns=['created_date'], values='deviceid')
ctab7t2=temp1.pivot_table(index=['Radio_station - #Device users(Tuned FM/Radio)'],columns=['created_date'], values='deviceid')
ctab7=ctab7t1.append(ctab7t2)

#Device Make
temp=data_process.groupby(['Dev_report','created_date']).agg({'deviceid':'nunique'}).reset_index()
temp.rename(columns={'Dev_report': 'Device make - #users'}, inplace=True)

temp1=data_process.groupby(['all','created_date']).agg({'deviceid':'nunique'}).reset_index()
temp1.rename(columns={'all': 'Device make - #users'}, inplace=True)

ctab8t1=temp.pivot_table(index=['Device make - #users'],columns=['created_date'], values='deviceid')
ctab8t2=temp1.pivot_table(index=['Device make - #users'],columns=['created_date'], values='deviceid')
ctab8=ctab8t1.append(ctab8t2)

                              

writer = pd.ExcelWriter('RawData_report.xlsx', engine='openpyxl')

ctab1.to_excel(writer, sheet_name='Sheet1', startrow=1)
ctab2.to_excel(writer, sheet_name='Sheet1', startrow=16)
ctab3.to_excel(writer, sheet_name='Sheet1', startrow=35)
ctab6.to_excel(writer, sheet_name='App_Radio', startrow=2)
ctab7.to_excel(writer, sheet_name='App_Radio', startrow=22)
ctab8.to_excel(writer, sheet_name='Device_Make', startrow=2)
ctab4.to_excel(writer, sheet_name='Album_name', startrow=2)
ctab5.to_excel(writer, sheet_name='Song_name', startrow=2)


writer.save()

def createFolder(directory):

#    if os.path.isdir("directory")==True:
#        return print ('Charts folder already exists')
    try:
        if not os.path.exists(directory):
            os.makedirs(directory)
    except OSError:
        print ('Error: Creating directory as already exists. ' +  directory)

createFolder('./Charts/')

#Create Charts
ax=plt.subplot()

#Event Name
temp=data_process.groupby(['keyname','created_date']).agg({'deviceid':'nunique'}).reset_index()
temp.rename(columns={'keyname': 'Event wise #Device users'}, inplace=True)
temp1=data_process.groupby(['all','created_date']).agg({'deviceid':'nunique'}).reset_index()
temp1.rename(columns={'all': 'Event wise #Device users'}, inplace=True)

#This portion not getting labels properly
#fig, ax = plt.subplots(figsize=(15,10))
#my_labels=temp['Event wise #Device users'].tolist()
#plt.legend(labels = my_labels)
#f_plot=temp.groupby('deviceid').plot(x='created_date',kind='line', y='Event wise #Device users', ax=ax, grid=True, legend=True, title='Number of device ids by eventname')

pivot_df = temp.pivot(index='created_date',columns='Event wise #Device users',values='deviceid')
pivot_df.to_csv('test.csv')
pivot_df=pd.read_csv('test.csv')
tick_labels = tuple(pivot_df['created_date'])

f_plot=pivot_df.plot.bar(stacked=True,figsize=(15,7), title='Number of users by event name', legend=True)
plt.xticks(range(0, len(pivot_df['created_date'])), tick_labels, rotation=45)
f_plot.legend(loc=2, bbox_to_anchor=(1, 1))
plt.xlabel('Date')
plt.ylabel('Number of devices')
fig = f_plot.get_figure()
fig.set_size_inches(15, 7)
fig.savefig('Charts/pngBarChart01.png', bbox_inches='tight', dpi=300)

pivot_df = temp.pivot(index='created_date',columns='Event wise #Device users',values='deviceid')
pivot_df.to_csv('test.csv')
pivot_df=pd.read_csv('test.csv')
tick_labels = tuple(pivot_df['created_date'])

f_plot=pivot_df.plot.area(stacked=True,figsize=(15,7), title='Number of users by event name', legend=True)
plt.xticks(range(0, len(pivot_df['created_date'])), tick_labels, rotation=45)
f_plot.legend(loc=2, bbox_to_anchor=(1, 1))
plt.xlabel('Date')
plt.ylabel('Number of devices')
fig = f_plot.get_figure()
fig.set_size_inches(15, 7)
fig.savefig('Charts/pngBarChart02.png', bbox_inches='tight', dpi=300)



#Mobile operator
temp=data_process.groupby(['mob_op_report','created_date']).agg({'deviceid':'nunique'}).reset_index()

pivot_df = temp.pivot(index='created_date',columns='mob_op_report',values='deviceid')
pivot_df.to_csv('test.csv')
pivot_df=pd.read_csv('test.csv')
tick_labels = tuple(pivot_df['created_date'])

f_plot=pivot_df.plot.bar(stacked=True,figsize=(15,7), title='Number of users by Mobile operators', legend=True)
plt.xticks(range(0, len(pivot_df['created_date'])), tick_labels, rotation=45)
f_plot.legend(loc=2, bbox_to_anchor=(1, 1))
plt.xlabel('Date')
plt.ylabel('Number of devices')
fig = f_plot.get_figure()
fig.set_size_inches(15, 7)
fig.savefig('Charts/pngBarChart03.png', bbox_inches='tight', dpi=300)


pivot_df = temp.pivot(index='created_date',columns='mob_op_report',values='deviceid')
pivot_df.to_csv('test.csv')
pivot_df=pd.read_csv('test.csv')
tick_labels = tuple(pivot_df['created_date'])

f_plot=pivot_df.plot.area(stacked=True,figsize=(15,7), title='Number of users by Mobile operators', legend=True)
plt.xticks(range(0, len(pivot_df['created_date'])), tick_labels, rotation=45)
f_plot.legend(loc=2, bbox_to_anchor=(1, 1))
plt.xlabel('Date')
plt.ylabel('Number of devices')
fig = f_plot.get_figure()
fig.set_size_inches(15, 7)
fig.savefig('Charts/pngBarChart04.png', bbox_inches='tight', dpi=300)

#Music App
temp=data_process[data_process['keyname']=='"Audio_Tuned"']
temp=temp.groupby(['app_rep','created_date']).agg({'deviceid':'nunique'}).reset_index()

pivot_df = temp.pivot(index='created_date',columns='app_rep',values='deviceid')
pivot_df.to_csv('test.csv')
pivot_df=pd.read_csv('test.csv')
tick_labels = tuple(pivot_df['created_date'])

f_plot=pivot_df.plot.bar(stacked=True,figsize=(15,7), title='Number of users by Music app', legend=True)
plt.xticks(range(0, len(pivot_df['created_date'])), tick_labels, rotation=45)
f_plot.legend(loc=2, bbox_to_anchor=(1, 1))
plt.xlabel('Date')
plt.ylabel('Number of devices')
fig = f_plot.get_figure()
fig.set_size_inches(15, 7)
fig.savefig('Charts/pngBarChart05.png', bbox_inches='tight', dpi=300)

pivot_df = temp.pivot(index='created_date',columns='app_rep',values='deviceid')
pivot_df.to_csv('test.csv')
pivot_df=pd.read_csv('test.csv')
tick_labels = tuple(pivot_df['created_date'])

f_plot=pivot_df.plot.area(stacked=True,figsize=(15,7), title='Number of users by Music app', legend=True)
plt.xticks(range(0, len(pivot_df['created_date'])), tick_labels, rotation=45)
f_plot.legend(loc=2, bbox_to_anchor=(1, 1))
plt.xlabel('Date')
plt.ylabel('Number of devices')
fig = f_plot.get_figure()
fig.set_size_inches(15, 7)
fig.savefig('Charts/pngBarChart06.png', bbox_inches='tight', dpi=300)


#Device make
temp=data_process.groupby(['Dev_report','created_date']).agg({'deviceid':'nunique'}).reset_index()

pivot_df = temp.pivot(index='created_date',columns='Dev_report',values='deviceid')
pivot_df.to_csv('test.csv')
pivot_df=pd.read_csv('test.csv')
tick_labels = tuple(pivot_df['created_date'])

f_plot=pivot_df.plot.bar(stacked=True,figsize=(15,7), title='Number of users by Device make', legend=True)
plt.xticks(range(0, len(pivot_df['created_date'])), tick_labels, rotation=45)
f_plot.legend(loc=2, bbox_to_anchor=(1, 1))
plt.xlabel('Date')
plt.ylabel('Number of devices')
fig = f_plot.get_figure()
fig.set_size_inches(15, 7)
fig.savefig('Charts/pngBarChart07.png', bbox_inches='tight', dpi=300)

pivot_df = temp.pivot(index='created_date',columns='Dev_report',values='deviceid')
pivot_df.to_csv('test.csv')
pivot_df=pd.read_csv('test.csv')
tick_labels = tuple(pivot_df['created_date'])

f_plot=pivot_df.plot.area(stacked=True,figsize=(15,7), title='Number of users by Device make', legend=True)
plt.xticks(range(0, len(pivot_df['created_date'])), tick_labels, rotation=45)
f_plot.legend(loc=2, bbox_to_anchor=(1, 1))
plt.xlabel('Date')
plt.ylabel('Number of devices')
fig = f_plot.get_figure()
fig.set_size_inches(15, 7)
fig.savefig('Charts/pngBarChart08.png', bbox_inches='tight', dpi=300)

#Radio channel
temp=data_process[data_process['keyname']=='"FM_Tuned"']
temp=temp.groupby(['Station Name','created_date']).agg({'deviceid':'nunique'}).reset_index()

pivot_df = temp.pivot(index='created_date',columns='Station Name',values='deviceid')
#pivot_df.columns = pivot_df.columns.droplevel().rename(None)
pivot_df=pd.read_csv('test.csv')

tick_labels = tuple(pivot_df['created_date'])
f_plot=pivot_df.plot.bar(stacked=True,figsize=(15,12), title='Number of users by FM stations', legend=True)
plt.xticks(range(0, len(pivot_df['created_date'])), tick_labels, rotation=45) 
f_plot.legend(loc=2, bbox_to_anchor=(1, 1))
plt.xlabel('Date')
plt.ylabel('Number of devices')
fig = f_plot.get_figure()
fig.set_size_inches(15, 7)
fig.savefig('Charts/pngBarChart09.png', bbox_inches='tight', dpi=300)

pivot_df = temp.pivot(index='created_date',columns='Station Name',values='deviceid')
pivot_df.to_csv('test.csv')
pivot_df=pd.read_csv('test.csv')
tick_labels = tuple(pivot_df['created_date'])
f_plot=pivot_df.plot.area(stacked=True,figsize=(15,12), title='Number of users by FM stations', legend=True)
plt.xticks(range(0, len(pivot_df['created_date'])), tick_labels, rotation=45) 
f_plot.legend(loc=2, bbox_to_anchor=(1, 1))
plt.xlabel('Date')
plt.ylabel('Number of devices')
fig = f_plot.get_figure()
fig.set_size_inches(15, 7)
fig.savefig('Charts/pngBarChart10.png', bbox_inches='tight', dpi=300)


#Alternative method to get x axis label in line graph
tick_labels = tuple(pivot_df['created_date'])
f_plot=pivot_df.plot(kind='line',stacked=False,figsize=(15,6))
plt.xticks(range(0, len(pivot_df['created_date'])), tick_labels, rotation=45) 
f_plot.legend(loc=2, bbox_to_anchor=(1, 1))
ax.set_xlabel('Date')
ax.set_ylabel('Number of devices')
# plt.show()
#Alternative method ends

os.remove('test.csv')

#Create PPT
OUTPUT_TAG = "KPI_Charts"
prs = pptx.Presentation()

slide = prs.slides.add_slide(prs.slide_layouts[0])

slide = prs.slides.add_slide(prs.slide_layouts[6])

# set title
title = slide.shapes.title
#title.text = OUTPUT_TAG

pic_left  = int(prs.slide_width * 0.05)
pic_top   = int(prs.slide_height * 0.1)
pic_width = int(prs.slide_width * 0.9)

for g in glob.glob("Charts/*.png"):
    print (g)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tb = slide.shapes.add_textbox(0, 0, prs.slide_width, pic_top / 2)
    p = tb.text_frame.add_paragraph()
    p.text = g
    p.font.size = pptx.util.Pt(14)

    img = scipy.misc.imread(g)
    pic_height = int(pic_width * img.shape[0] / img.shape[1])
    #pic   = slide.shapes.add_picture(g, pic_left, pic_top)
    pic   = slide.shapes.add_picture(g, pic_left, pic_top, pic_width, pic_height)

prs.save("Charts/%s.pptx" % OUTPUT_TAG)

